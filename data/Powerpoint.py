from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "State Quality & Expected Utility Pipeline"
subtitle = slide.placeholders[1]
subtitle.text = "CS 5260 Programming Project 1"

# Slide 2: State Quality Function
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "State Quality Function"
body = slide.shapes.placeholders[1].text_frame
body.text = "We define the State Quality function Q(c) as:"
p = body.add_paragraph()
p.text = "Q(c) = (1 / Population) × Σ_r [ w_r × (A_r − b_r × Population) ]"
p.level = 1
p = body.add_paragraph()
p.text = "• w_r: weight for resource r"
p.level = 2
p = body.add_paragraph()
p.text = "• b_r: per-capita baseline for resource r"
p.level = 2

# Slide 3: Justification
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Justification for State Quality"
body = slide.shapes.placeholders[1].text_frame
body.text = "• Weighted-sum multi-criteria analysis is a standard approach" 
p = body.add_paragraph()
p.text = "  (GeeksforGeeks, 2025) [1]"
p.level = 1
p = body.add_paragraph()
p.text = "• Normalization by population mirrors Human Development Index (UNDP, 2024) [2]"
p.level = 1
p = body.add_paragraph()
p.text = "• Waste terms included with negative weights to penalize pollution"
p.level = 1

# Slide 4: Score Pipeline Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Pipeline of Score Calculations"
body = slide.shapes.placeholders[1].text_frame
body.text = "1. Compute Q_start and Q_end"
p = body.add_paragraph()
p.text = "2. Undiscounted Reward: R = Q_end − Q_start"
p.level = 1
p = body.add_paragraph()
p.text = "3. Discounted Reward: DR = γ^N × R"
p.level = 1
p = body.add_paragraph()
p.text = "4. Acceptance Probability: P = logistic_k(DR − x0)"
p.level = 1
p = body.add_paragraph()
p.text = "5. Expected Utility: EU = P×DR + (1−P)×C"
p.level = 1

# Slide 5: Parameter Settings & Insights
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Parameter Settings & Insights"
body = slide.shapes.placeholders[1].text_frame
body.text = "• Discount factor γ = 0.9: balances near vs. future rewards"
p = body.add_paragraph()
p.text = "• Logistic steepness k = 1.0, midpoint x0 = 0: neutral accept"
p.level = 1
p = body.add_paragraph()
p.text = "• Failure cost C = -10: moderate penalty for unaccepted schedules"
p.level = 1
p = body.add_paragraph()
p.text = "• Depth = 6, Beam width = 50: initial exploration limits"
p.level = 1

# Save the presentation
output_path = "/mnt/data/ExplanationSlides.pptx"
prs.save(output_path)

output_path
